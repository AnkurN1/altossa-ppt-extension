from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image
import requests
from io import BytesIO
from urllib.parse import urlparse

# Keywords to detect product type
PRODUCT_TYPES = ["sofa", "chair", "bed", "table", "armchair", "desk", "bench", "stool", "cabinet"]

def extract_product_type(url: str) -> str:
    url_lower = url.lower()
    for type_ in PRODUCT_TYPES:
        if type_ in url_lower:
            return type_.capitalize()
    return "Product"

def shorten_url(url: str) -> str:
    parsed = urlparse(url)
    return f"{parsed.netloc}{parsed.path}".strip("/")

def generate_ppt(data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    for item in data:
        slide = prs.slides.add_slide(blank_slide_layout)
        img_url = item.get('image')
        link = item.get('page', '')
        product_type = extract_product_type(link)
        short_link = shorten_url(link)

        try:
            response = requests.get(img_url)
            response.raise_for_status()
            img_bytes = BytesIO(response.content)

            # Convert non-supported formats to JPEG via PIL
            image = Image.open(img_bytes).convert("RGB")
            jpeg_stream = BytesIO()
            image.save(jpeg_stream, format="JPEG")
            jpeg_stream.seek(0)

            # Define layout zone
            top_title_space = Inches(0.7)
            bottom_space = Inches(0.5)
            
            usable_width = prs.slide_width
            usable_height = prs.slide_height - top_title_space - bottom_space
            
            # Get image size
            img_width_px, img_height_px = image.size
            
            # Scale proportionally
            width_ratio = usable_width / img_width_px
            height_ratio = usable_height / img_height_px
            scale_ratio = min(width_ratio, height_ratio)
            
            final_width = img_width_px * scale_ratio
            final_height = img_height_px * scale_ratio
            
            # Left aligned, vertically centered
            left = Inches(0.2)
            top = top_title_space + (usable_height - final_height) / 2
            
            slide.shapes.add_picture(
                jpeg_stream,
                left,
                top,
                width=final_width,
                height=final_height
            )

            # Top-left: Product type label
            title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(5), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = product_type
            title_frame.paragraphs[0].font.size = Pt(16)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

            # Bottom-left: Shortened link
            link_box = slide.shapes.add_textbox(Inches(0.15), Inches(7.2), Inches(6), Inches(0.5))
            link_frame = link_box.text_frame
            link_frame.text = f"{short_link}"
            link_frame.paragraphs[0].font.size = Pt(10)
            link_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)


        except Exception as e:
            print("❌ Failed to process image:", img_url)
            print("Error:", e)

    prs.save("output.pptx")
    return "output.pptx"
